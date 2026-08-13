"""Builds the technical walkthrough deck: how the demo works and every number behind it.

Companion to `create_presentation.py` (the executive/investor deck). This one is aimed at
engineers and reviewers who need to see the actual equations, coefficients and thresholds,
so every figure here is transcribed straight from `src/lib/damModel.ts`. If that file
changes, re-run this script.

    python create_technical_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUT = "Demo_Dam_How_It_Works.pptx"

# Palette lifted from the dashboard shell so the deck and the app read as one product.
BG = RGBColor(11, 15, 25)          # #0b0f19
CARD = RGBColor(22, 29, 45)        # #161d2d
CARD_DEEP = RGBColor(15, 22, 36)
CYAN = RGBColor(0, 229, 255)
PURPLE = RGBColor(139, 92, 246)
EMERALD = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
PINK = RGBColor(232, 74, 156)
TEXT = RGBColor(243, 244, 246)
MUTED = RGBColor(156, 163, 175)
LINE = RGBColor(55, 65, 81)

BODY_FONT = "Arial"
MONO_FONT = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------
def new_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def card(slide, left, top, width, height, fill=CARD, border=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    if radius:
        try:  # tighter corner radius than the PowerPoint default
            shape.adjustments[0] = 0.08
        except (IndexError, KeyError):
            pass
    return shape


def text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = list of (string, size_pt, bold, color, font_name_or_None, space_before_pt)."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (s, size, bold, color, font, space) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = s
        p.alignment = align
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font or BODY_FONT
        if space:
            p.space_before = Pt(space)
    return box


def header(slide, title, subtitle=None, accent=CYAN):
    card(slide, 0.0, 0.0, 0.14, 7.5, fill=accent, radius=False)
    runs = [(title, 25, True, accent, None, 0)]
    if subtitle:
        runs.append((subtitle, 12.5, False, MUTED, None, 5))
    text(slide, 0.62, 0.42, 12.1, 1.0, runs)


def footnote(slide, note):
    text(slide, 0.62, 6.92, 12.1, 0.35, [(note, 9.5, False, RGBColor(100, 110, 128), None, 0)])


def formula(slide, left, top, width, lines, accent=CYAN, height=None):
    h = height if height else 0.42 + 0.32 * len(lines)
    card(slide, left, top, width, h, fill=CARD_DEEP, border=accent)
    runs = [(ln, 14, i == 0, TEXT if i == 0 else MUTED, MONO_FONT, 0 if i == 0 else 6)
            for i, ln in enumerate(lines)]
    text(slide, left + 0.28, top + 0.2, width - 0.56, h - 0.3, runs)
    return top + h


def grid(slide, left, top, width, headers, rows, col_w, row_h=0.38, accent=CYAN,
         row_colors=None, mono_cols=(), size=10.5):
    """Hand-drawn table: rectangles + textboxes, so every cell can carry its own colour."""
    scale = width / sum(col_w)
    widths = [c * scale for c in col_w]

    head = card(slide, left, top, width, row_h, fill=RGBColor(30, 41, 59), radius=False)
    head.line.color.rgb = accent
    head.line.width = Pt(1)
    x = left
    for w, h_text in zip(widths, headers):
        text(slide, x + 0.12, top + 0.09, w - 0.2, row_h,
             [(h_text, size, True, accent, None, 0)])
        x += w

    y = top + row_h
    for r, row in enumerate(rows):
        band = CARD if r % 2 == 0 else CARD_DEEP
        strip = card(slide, left, y, width, row_h, fill=band, radius=False)
        strip.line.color.rgb = LINE
        strip.line.width = Pt(0.5)
        x = left
        for c, (w, cell) in enumerate(zip(widths, row)):
            colour = TEXT
            if row_colors and row_colors[r] and c == len(row) - 1:
                colour = row_colors[r]
            elif c == 0:
                colour = CYAN
            text(slide, x + 0.12, y + 0.09, w - 0.2, row_h,
                 [(str(cell), size, c == 0, colour, MONO_FONT if c in mono_cols else None, 0)])
            x += w
        y += row_h
    return y


def bullets(slide, left, top, width, items, accent=CYAN, size=12, gap=0.86, card_h=0.74):
    y = top
    for title, body in items:
        card(slide, left, y, width, card_h, fill=CARD, border=None)
        text(slide, left + 0.24, y + 0.1, width - 0.48, card_h - 0.16,
             [(title, size + 1.5, True, accent, None, 0),
              (body, size - 0.5, False, MUTED, None, 3)])
        y += gap
    return y


def chip(slide, left, top, width, height, value, label, colour):
    card(slide, left, top, width, height, fill=CARD_DEEP, border=colour)
    text(slide, left, top + 0.16, width, height,
         [(value, 21, True, colour, MONO_FONT, 0), (label, 10.5, False, MUTED, None, 4)],
         align=PP_ALIGN.CENTER)


def arrow(slide, left, top, width, height, colour=CYAN):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    a.fill.solid()
    a.fill.fore_color.rgb = colour
    a.line.fill.background()
    a.shadow.inherit = False
    return a


# ---------------------------------------------------------------------------
# 1 — Title
# ---------------------------------------------------------------------------
s = new_slide()
card(s, 0.8, 1.05, 11.733, 4.15, fill=CARD, border=CYAN)
text(s, 1.25, 1.45, 10.8, 2.6, [
    ("DEMO DAM · DIGITAL TWIN", 13, True, MUTED, MONO_FONT, 0),
    ("How It Works", 42, True, CYAN, None, 6),
    ("The model, the equations, and every constant behind the dashboard", 19, False, TEXT, None, 10),
    ("A technical walkthrough of the simulation layer: three operator inputs drive six "
     "sensor responses, one risk decomposition, and three scripted scenarios — all as "
     "pure, deterministic functions.", 12.5, False, MUTED, None, 14),
])
labels = [("3", "Sim inputs", CYAN), ("6", "Sensor nodes", PURPLE),
          ("2", "Core equations", EMERALD), ("3", "Scenarios", AMBER)]
for i, (v, l, c) in enumerate(labels):
    chip(s, 1.25 + i * 2.62, 3.85, 2.35, 1.05, v, l, c)
text(s, 0.8, 5.45, 11.733, 0.6, [
    ("SU YAPI Engineering & Consulting   ·   Para Sidara", 12, True, TEXT, None, 0),
    ("Demonstration dataset — figures are simulated, not live field telemetry.",
     10.5, False, MUTED, None, 4),
])

# ---------------------------------------------------------------------------
# 2 — Reading map
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "What this deck covers", "Four questions, answered in order. Skip to the one you need.")
bullets(s, 0.62, 1.55, 12.1, [
    ("1 · How does data flow?",
     "Three slider inputs → pure functions → sensors, risk, protocol, 3D twin. No hidden state, no server."),
    ("2 · How is each sensor reading calculated?",
     "One linear superposition equation with per-sensor sensitivity coefficients. Slides 5–9."),
    ("3 · How is the risk score calculated?",
     "Three weighted drivers plus a baseline, clamped and banded into three status levels. Slides 10–12."),
    ("4 · How do the scripted scenarios run?",
     "Keyframe interpolation on a smoothstep curve, with events fired at progress thresholds. Slides 13–15."),
    ("5 · What is this model NOT?",
     "The honest limits — what a reviewer should know before reading anything into the numbers. Slide 17."),
], gap=1.0, card_h=0.86)
footnote(s, "Every constant in this deck is transcribed from src/lib/damModel.ts.")

# ---------------------------------------------------------------------------
# 3 — Architecture / data flow
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "1 · Data flow", "One direction, top to bottom. Everything downstream is derived on render.")

card(s, 0.62, 1.5, 3.5, 1.5, fill=CARD, border=CYAN)
text(s, 0.82, 1.66, 3.1, 1.2, [
    ("INPUTS", 11, True, CYAN, MONO_FONT, 0),
    ("surge · pga · spillwayGate", 12, False, TEXT, MONO_FONT, 6),
    ("Set by sliders, or driven by a running scenario.", 10, False, MUTED, None, 5),
])
arrow(s, 4.32, 2.05, 0.7, 0.4)

card(s, 5.22, 1.5, 3.5, 1.5, fill=CARD, border=PURPLE)
text(s, 5.42, 1.66, 3.1, 1.2, [
    ("PURE FUNCTIONS", 11, True, PURPLE, MONO_FONT, 0),
    ("readSensors()\ncomputeRisk()", 12, False, TEXT, MONO_FONT, 6),
    ("No effects, no clock, no randomness.", 10, False, MUTED, None, 5),
])
arrow(s, 8.92, 2.05, 0.7, 0.4, PURPLE)

card(s, 9.82, 1.5, 2.9, 1.5, fill=CARD, border=EMERALD)
text(s, 10.02, 1.66, 2.5, 1.2, [
    ("DERIVED UI", 11, True, EMERALD, MONO_FONT, 0),
    ("6 readings + risk\n+ protocol text", 12, False, TEXT, MONO_FONT, 6),
    ("Recomputed every render.", 10, False, MUTED, None, 5),
])

text(s, 0.62, 3.25, 12.1, 0.4,
     [("The four consumers of that derived state", 13, True, TEXT, None, 0)])
consumers = [
    ("3D twin", "Sensor beacons colour by status; labels show live values.", CYAN),
    ("Telemetry charts", "28-point rolling history, sampled on a 2 s clock.", PURPLE),
    ("Risk gauge", "Score plus the three named driver contributions.", PINK),
    ("Protocol panel", "Recommendation text selected by risk band.", EMERALD),
]
for i, (t_, d, c) in enumerate(consumers):
    x = 0.62 + i * 3.06
    card(s, x, 3.72, 2.86, 1.35, fill=CARD, border=c)
    text(s, x + 0.2, 3.9, 2.46, 1.1,
         [(t_, 13, True, c, None, 0), (d, 10.5, False, MUTED, None, 5)])

formula(s, 0.62, 5.35, 12.1, [
    "Why it matters: state is never synced, only derived.",
    "The same three numbers always produce the same six readings and the same score —",
    "so any frame of the demo is exactly reproducible from three values.",
])

# ---------------------------------------------------------------------------
# 4 — The three inputs
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "2 · The three simulator inputs", "Everything the demo shows is a function of these three numbers.")
grid(s, 0.62, 1.55, 12.1,
     ["Input", "Meaning", "Range", "Step", "Rest value"],
     [
         ["surge", "Reservoir rise above normal operating level", "0 → 15 m", "0.5", "0"],
         ["pga", "Peak ground acceleration at the abutment", "0 → 0.45 g", "0.01", "0.02"],
         ["spillwayGate", "Gate opening — discharge capacity available", "10 → 100 %", "5", "100"],
     ],
     col_w=[2.0, 5.0, 1.9, 1.0, 1.4], row_h=0.46, mono_cols=(0, 2, 3, 4), size=11.5)

formula(s, 0.62, 3.35, 12.1, [
    "closure = 100 - spillwayGate",
    "Gate opening is inverted into a closure percentage before it reaches any equation.",
    "A fully open gate (100%) contributes exactly zero. A seized gate at 10% contributes 90 units of closure.",
], accent=EMERALD)

bullets(s, 0.62, 5.05, 12.1, [
    ("Rest state is not zero risk",
     "PGA idles at 0.02 g (ambient microseismic background), so the dashboard never shows a dead-flat trace."),
    ("Sliders lock during a scenario",
     "While a scripted scenario runs, the scenario clock owns the three inputs and the sliders go read-only."),
], gap=0.82, card_h=0.7)

# ---------------------------------------------------------------------------
# 5 — Sensor network
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "3 · The sensor network", "Six instruments, each with a rest reading and two alarm thresholds.")
grid(s, 0.62, 1.5, 12.1,
     ["ID", "Instrument", "Location", "Unit", "Base", "Warn ≥", "Crit ≥"],
     [
         ["P-01", "Piezometer", "Upstream Heel (EL 490 m)", "kPa", "242.4", "270", "300"],
         ["P-02", "Piezometer", "Foundation Gallery (EL 510 m)", "kPa", "185.1", "205", "220"],
         ["P-03", "Piezometer", "Downstream Toe (EL 480 m)", "kPa", "92.6", "130", "150"],
         ["INC-04", "Pendulum / Inclinometer", "Crest Centre (Block 04)", "mm", "3.2", "6.5", "8.0"],
         ["SF-02", "Seepage Flow Gauge", "Drainage Gallery G-02", "L/min", "14.2", "20", "25"],
         ["WL-01", "Water Level Radar", "Reservoir Intake Tower", "m", "537.4", "543", "546"],
     ],
     col_w=[1.1, 2.6, 3.5, 0.9, 1.1, 1.0, 1.0], row_h=0.44, mono_cols=(0, 3, 4, 5, 6), size=11)

text(s, 0.62, 4.55, 12.1, 0.4,
     [("Base = the reading at rest: surge 0, pga 0, gate 100%. Thresholds are one-sided — "
       "every sensor here alarms on rising values only.", 12, False, MUTED, None, 0)])

formula(s, 0.62, 5.2, 12.1, [
    "Instrument naming follows arch-dam practice, not embankment practice.",
    "Heel / toe / foundation gallery / pendulum — the terminology a concrete double-curvature",
    "arch dam would actually use on its instrumentation schedule.",
], accent=PURPLE)

# ---------------------------------------------------------------------------
# 6 — Core equation
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "4 · The sensor response equation", "One line of arithmetic produces every reading on screen.")

card(s, 0.62, 1.5, 12.1, 1.35, fill=CARD_DEEP, border=CYAN)
text(s, 0.62, 1.72, 12.1, 1.0, [
    ("value  =  base  +  surge × kSurge  +  pga × kPga  +  closure × kGate",
     19, True, CYAN, MONO_FONT, 0),
    ("linear superposition of three independent load cases", 12, False, MUTED, None, 8),
], align=PP_ALIGN.CENTER)

terms = [
    ("base", "Rest reading", "The instrument's value with the reservoir at normal level and no shaking.", TEXT),
    ("surge × kSurge", "Hydrostatic term", "Reservoir rise in metres × that sensor's metres-of-surge sensitivity.", CYAN),
    ("pga × kPga", "Seismic term", "Ground acceleration in g × that sensor's per-g sensitivity.", PURPLE),
    ("closure × kGate", "Discharge term", "Percent of gate closure × that sensor's per-percent sensitivity.", EMERALD),
]
y = 3.1
for expr, name, desc, colour in terms:
    card(s, 0.62, y, 12.1, 0.82, fill=CARD)
    text(s, 0.86, y + 0.13, 3.2, 0.6, [(expr, 13, True, colour, MONO_FONT, 0)])
    text(s, 4.2, y + 0.06, 8.3, 0.72,
         [(name, 12.5, True, TEXT, None, 0), (desc, 10.5, False, MUTED, None, 3)])
    y += 0.94

formula(s, 0.62, 6.15, 12.1, [
    "Then: value is rounded to the sensor's own decimal precision before display and comparison.",
], accent=AMBER, height=0.68)

# ---------------------------------------------------------------------------
# 7 — Coefficient matrix
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "5 · The sensitivity coefficients", "Which loads each instrument actually responds to — and how strongly.")
grid(s, 0.62, 1.5, 12.1,
     ["Sensor", "kSurge (per m)", "kPga (per g)", "kGate (per %)", "Dominant driver"],
     [
         ["P-01", "8.5", "150", "0.15", "Hydrostatic"],
         ["P-02", "5.2", "90", "0.10", "Hydrostatic"],
         ["P-03", "2.4", "60", "0.08", "Mixed"],
         ["INC-04", "0.4", "25", "0", "Seismic"],
         ["SF-02", "1.8", "30", "0", "Hydrostatic"],
         ["WL-01", "1.0", "0", "0.02", "Hydrostatic"],
     ],
     col_w=[1.6, 2.5, 2.3, 2.4, 3.3], row_h=0.44, mono_cols=(0, 1, 2, 3), size=11,
     row_colors=[CYAN, CYAN, AMBER, PURPLE, CYAN, CYAN])

reads = [
    ("A zero means genuinely no coupling", PURPLE,
     "INC-04 and SF-02 carry kGate = 0: crest deflection and gallery seepage do not respond to "
     "gate position in this model. WL-01 carries kPga = 0 — shaking does not move the reservoir surface."),
    ("WL-01's kSurge is exactly 1.0", CYAN,
     "By construction: reservoir level rises one metre per metre of surge. It is the one coefficient "
     "that is definitional rather than calibrated."),
    ("Piezometers dominate the seismic term", AMBER,
     "kPga of 150 at the heel means 0.45 g alone adds 67.5 kPa — more than half the distance from base "
     "to the critical limit."),
]
y = 4.6
for title, colour, body in reads:
    card(s, 0.62, y, 12.1, 0.72, fill=CARD)
    text(s, 0.86, y + 0.1, 11.6, 0.6,
         [(title, 12.5, True, colour, None, 0), (body, 10.5, False, MUTED, None, 3)])
    y += 0.8

# ---------------------------------------------------------------------------
# 8 — Worked example: sensors
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "6 · Worked example — all six readings",
       "Operating point: surge = 4.0 m, pga = 0.05 g, gate = 100% (closure = 0).", accent=AMBER)

formula(s, 0.62, 1.45, 12.1, [
    "P-01  =  242.4  +  4.0 × 8.5  +  0.05 × 150  +  0 × 0.15",
    "      =  242.4  +  34.0  +  7.5  +  0.0   =   283.9 kPa      →  ≥ 270 warn, < 300 crit  →  WARNING",
], accent=AMBER, height=1.05)

grid(s, 0.62, 2.75, 12.1,
     ["Sensor", "Base", "+ Surge", "+ PGA", "+ Gate", "Value", "Status"],
     [
         ["P-01", "242.4", "34.0", "7.5", "0.0", "283.9 kPa", "WARNING"],
         ["P-02", "185.1", "20.8", "4.5", "0.0", "210.4 kPa", "WARNING"],
         ["P-03", "92.6", "9.6", "3.0", "0.0", "105.2 kPa", "NORMAL"],
         ["INC-04", "3.2", "1.6", "1.25", "0.0", "6.1 mm", "NORMAL"],
         ["SF-02", "14.2", "7.2", "1.5", "0.0", "22.9 L/min", "WARNING"],
         ["WL-01", "537.4", "4.0", "0.0", "0.0", "541.4 m", "NORMAL"],
     ],
     col_w=[1.5, 1.5, 1.5, 1.5, 1.5, 2.3, 2.0], row_h=0.44, mono_cols=(0, 1, 2, 3, 4, 5), size=11,
     row_colors=[AMBER, AMBER, EMERALD, EMERALD, AMBER, EMERALD])

text(s, 0.62, 5.85, 12.1, 0.8, [
    ("Three sensors in warning, three normal — from a 4 m surge and near-negligible shaking. "
     "The spread is the point: because each instrument carries its own coefficients and its own "
     "thresholds, one operating point produces a mixed picture rather than six identical bars.",
     12, False, MUTED, None, 0),
])
footnote(s, "INC-04 sums to 6.05 and displays as 6.1 — rounding to the sensor's own precision happens before display and before the threshold test.")

# ---------------------------------------------------------------------------
# 9 — Status + utilisation
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "7 · Status bands and bar fill", "How a number becomes a colour.")

formula(s, 0.62, 1.5, 5.85, [
    "status =",
    "  value ≥ critAt  →  'critical'",
    "  value ≥ warnAt  →  'warning'",
    "  otherwise       →  'normal'",
], accent=CYAN, height=1.75)

formula(s, 6.87, 1.5, 5.85, [
    "utilisation =",
    "  min(1.4, value / critAt)",
    "",
    "Drives the fill fraction on each bar.",
], accent=PURPLE, height=1.75)

text(s, 0.62, 3.5, 12.1, 0.4,
     [("Two details worth knowing", 13.5, True, TEXT, None, 0)])

notes = [
    ("Comparison happens after rounding", CYAN,
     "The value is rounded to the sensor's display precision first, then tested. What you read on "
     "screen is exactly what was compared — a reading shown as 270.0 kPa is in warning, with no "
     "hidden 269.97 behind it."),
    ("Utilisation is capped at 1.4, not 1.0", PURPLE,
     "Bars keep growing past the critical limit so an overshoot stays visible, but stop at 140% so a "
     "runaway value cannot blow out the layout. A bar pinned at full width means ≥ 140% of critical."),
]
y = 3.95
for title, colour, body in notes:
    card(s, 0.62, y, 12.1, 1.0, fill=CARD, border=colour)
    text(s, 0.9, y + 0.14, 11.5, 0.85,
         [(title, 13, True, colour, None, 0), (body, 11, False, MUTED, None, 4)])
    y += 1.15

formula(s, 0.62, 6.3, 12.1, [
    "Status drives colour everywhere at once: the sensor tile, its 3D beacon, and its label ring.",
], accent=EMERALD, height=0.62)

# ---------------------------------------------------------------------------
# 10 — Risk equation
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "8 · The risk score", "A sum of named contributions, so the gauge can explain why it moved.")

card(s, 0.62, 1.45, 12.1, 1.15, fill=CARD_DEEP, border=PINK)
text(s, 0.62, 1.68, 12.1, 0.8, [
    ("score  =  clamp( hydrostatic + seismic + discharge + 12 ,  5 ,  100 )",
     18, True, PINK, MONO_FONT, 0),
], align=PP_ALIGN.CENTER)

drivers = [
    ("hydrostatic", "surge × 3.8", "0 → 57.0", "Reservoir loading. The largest single contributor at full range.", CYAN),
    ("seismic", "pga × 140", "0 → 63.0", "Ground shaking. Steepest coefficient — 0.45 g alone exceeds the warning band.", PURPLE),
    ("discharge", "closure × 0.35", "0 → 31.5", "Lost spillway capacity. Smallest contributor, but always present when gates are throttled.", PINK),
    ("baseline", "12 (constant)", "12", "RISK_BASELINE — the standing risk of operating a dam at all. Never zero.", MUTED),
]
y = 2.85
for name, expr, rng, desc, colour in drivers:
    card(s, 0.62, y, 12.1, 0.78, fill=CARD)
    text(s, 0.86, y + 0.19, 2.1, 0.5, [(name, 12.5, True, colour, MONO_FONT, 0)])
    text(s, 3.0, y + 0.19, 2.0, 0.5, [(expr, 12, False, TEXT, MONO_FONT, 0)])
    text(s, 5.1, y + 0.19, 1.5, 0.5, [(rng, 11.5, False, MUTED, MONO_FONT, 0)])
    text(s, 6.8, y + 0.15, 5.7, 0.6, [(desc, 10.5, False, MUTED, None, 0)])
    y += 0.86

formula(s, 0.62, 6.35, 12.1, [
    "Sum at full range = 57 + 63 + 31.5 + 12 = 163.5, so the score saturates at 100 well before the sliders max out.",
], accent=RED, height=0.62)

# ---------------------------------------------------------------------------
# 11 — Risk bands + worked example
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "9 · Risk bands and a worked score", "Three bands, and the same operating point as slide 6.", accent=AMBER)

bands = [("score > 60", "CRITICAL RISK", "SEVERE STRESS", RED),
         ("score > 28", "ELEVATED RISK", "WARNING", AMBER),
         ("score ≤ 28", "LOW RISK", "STRUCTURALLY STABLE", EMERALD)]
for i, (cond, label, sub, colour) in enumerate(bands):
    x = 0.62 + i * 4.08
    card(s, x, 1.5, 3.86, 1.35, fill=CARD_DEEP, border=colour)
    text(s, x + 0.24, 1.68, 3.4, 1.1,
         [(cond, 13, True, colour, MONO_FONT, 0),
          (label, 14, True, TEXT, None, 6),
          (sub, 10.5, False, MUTED, None, 3)])

formula(s, 0.62, 3.15, 12.1, [
    "At surge = 4.0 m, pga = 0.05 g, gate = 100%:",
    "hydrostatic = 4.0 × 3.8   = 15.2",
    "seismic     = 0.05 × 140  =  7.0",
    "discharge   = 0 × 0.35    =  0.0",
    "baseline                  = 12.0",
    "score = 15.2 + 7.0 + 0.0 + 12.0 = 34.2   →  > 28  →  ELEVATED RISK — WARNING",
], accent=AMBER, height=2.35)

text(s, 0.62, 5.75, 12.1, 1.0, [
    ("The gauge shows the split, not just the total.", 13, True, TEXT, None, 0),
    ("34.2 points made of 15.2 hydrostatic + 7.0 seismic + 0.0 discharge tells an operator to look "
     "at the reservoir, not the gates or the fault. That decomposition is the whole reason the score "
     "is a sum of named terms rather than a single opaque formula.", 11.5, False, MUTED, None, 5),
])

# ---------------------------------------------------------------------------
# 12 — Protocol engine
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "10 · The recommended protocol", "A lookup on the risk band — deterministic, not generative.")

protocols = [
    ("CRITICAL", RED, "score > 60",
     "Open spillway gates to 100%. Dispatch structural team to inspect the downstream toe (P-03) and "
     "drainage gallery. Notify downstream settlements per Emergency Action Plan."),
    ("WARNING", AMBER, "score > 28",
     "Hydrodynamic pressure elevated. Increase piezometric sampling to 1 s intervals, verify gallery "
     "drain flow at SF-02, and stage auxiliary spillway gates for release."),
    ("NORMAL", EMERALD, "score ≤ 28",
     "Structure stable under nominal hydrodynamic loads. No intervention required. Piezometric heads "
     "and crest deflection within design safety envelopes."),
]
y = 1.5
for level, colour, cond, body in protocols:
    card(s, 0.62, y, 12.1, 1.42, fill=CARD, border=colour)
    text(s, 0.9, y + 0.18, 2.3, 1.0,
         [(level, 15, True, colour, None, 0), (cond, 11, False, MUTED, MONO_FONT, 5)])
    text(s, 3.35, y + 0.2, 9.1, 1.05, [(body, 12, False, TEXT, None, 0)])
    y += 1.58

formula(s, 0.62, 6.3, 12.1, [
    "No language model is involved. The 'AI Recommended Protocol' panel selects one of three fixed strings.",
], accent=CYAN, height=0.62)

# ---------------------------------------------------------------------------
# 13 — Scenario engine
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "11 · The scenario engine", "How a scripted event drives the sliders, the twin, the charts and the gauge from one clock.")

steps = [
    ("Progress", "p = elapsed / durationMs, clamped to 1. One requestAnimationFrame loop owns it."),
    ("Bracket", "Find the keyframe pair a, b where a.at ≤ p ≤ b.at."),
    ("Normalise", "raw = (p − a.at) / (b.at − a.at) — position between those two keyframes."),
    ("Ease", "t = raw² × (3 − 2·raw) — smoothstep, so slider motion accelerates and settles."),
    ("Interpolate", "Each of the three inputs is lerped from a to b using t, then pushed to state."),
]
y = 1.5
for i, (name, body) in enumerate(steps):
    card(s, 0.62, y, 12.1, 0.72, fill=CARD)
    text(s, 0.86, y + 0.18, 0.5, 0.5, [(str(i + 1), 15, True, CYAN, MONO_FONT, 0)])
    text(s, 1.5, y + 0.08, 11.0, 0.62,
         [(name, 12.5, True, TEXT, None, 0), (body, 11, False, MUTED, MONO_FONT, 3)])
    y += 0.8

formula(s, 0.62, 5.6, 5.85, [
    "Events fire on threshold crossing:",
    "while (fired < events.length",
    "       && events[fired].at <= p)",
    "  → append to the log, once only",
], accent=PURPLE, height=1.55)

formula(s, 6.87, 5.6, 5.85, [
    "Why smoothstep, not linear:",
    "linear keyframes make the sliders",
    "snap between poses. Smoothstep eases",
    "in and out — it reads as real motion.",
], accent=EMERALD, height=1.55)

# ---------------------------------------------------------------------------
# 14 — Scenario catalogue
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "12 · The three scenarios", "Each one exercises a different term of the risk equation.")
grid(s, 0.62, 1.5, 12.1,
     ["Scenario", "Premise", "Runtime", "Peak load", "Camera", "Events"],
     [
         ["Flood Event", "1-in-500-year inflow, staged gate response", "26 s", "surge 15.0 m", "WL-01", "8"],
         ["Seismic Event", "M6.4 near-field rupture, aftershock sequence", "22 s", "pga 0.38 g", "INC-04", "8"],
         ["Gate Malfunction", "Radial gate seizure during high inflow", "20 s", "gate 10%", "SF-02", "7"],
     ],
     col_w=[2.3, 4.6, 1.2, 1.9, 1.2, 0.9], row_h=0.5, mono_cols=(2, 3, 4, 5), size=11)

text(s, 0.62, 3.35, 12.1, 0.4,
     [("What each one is designed to show", 13.5, True, TEXT, None, 0)])

shows = [
    ("Flood", CYAN, "Drives the hydrostatic term to its ceiling. Gates stage down 100 → 75 → 45 → 30% as "
                    "capacity saturates, then reopen fully at p = 0.88 for drawdown. Ends with the structure recovered."),
    ("Seismic", PURPLE, "Spikes the seismic term to 0.38 g at p = 0.12, decays, then a second 0.29 g aftershock "
                        "at p = 0.42. Surge barely moves — this scenario isolates INC-04 crest deflection."),
    ("Gate", PINK, "Holds pga flat and collapses the gate to 10%, so the discharge term and a slow surge climb "
                   "carry the whole score. Recovery only comes at p = 0.86 when the gate is freed."),
]
y = 3.8
for name, colour, body in shows:
    card(s, 0.62, y, 12.1, 0.86, fill=CARD)
    text(s, 0.88, y + 0.24, 1.5, 0.5, [(name, 13, True, colour, None, 0)])
    text(s, 2.4, y + 0.12, 10.1, 0.72, [(body, 11, False, MUTED, None, 0)])
    y += 0.94

formula(s, 0.62, 6.65, 12.1, [
    "All three end in a recovered state — the demo is designed to close on a stable structure, not a failure.",
], accent=EMERALD, height=0.62)

# ---------------------------------------------------------------------------
# 15 — Telemetry sampling
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "13 · Telemetry sampling and the trace jitter",
       "The one place the demo is deliberately non-deterministic.")

for i, (v, l, c) in enumerate([("2 s", "Sample interval", CYAN), ("28", "Points retained", PURPLE),
                               ("56 s", "Visible window", EMERALD), ("±0.075%", "Jitter amplitude", AMBER)]):
    chip(s, 0.62 + i * 3.06, 1.5, 2.86, 1.15, v, l, c)

formula(s, 0.62, 2.95, 12.1, [
    "jitter = (random() − 0.5) × max(0.05, value × 0.0015)",
    "sample = round(value + jitter, decimals)",
], accent=AMBER, height=1.15)

notes = [
    ("Why add noise at all", CYAN,
     "A pure model produces a perfectly flat line when nothing moves. Real instruments never do. "
     "The jitter is cosmetic — it makes an idle trace look like an instrument rather than a constant."),
    ("It never changes a status", EMERALD,
     "Jitter is applied only to the charted history. Sensor tiles, statuses, the risk score and the "
     "3D twin all read the clean deterministic value. Noise can never trip an alarm."),
    ("The floor matters", PURPLE,
     "max(0.05, …) keeps a visible wiggle on small readings — INC-04 at 3.2 mm would otherwise jitter "
     "by ±0.0024 mm and round away to a flat line."),
]
y = 4.35
for title, colour, body in notes:
    card(s, 0.62, y, 12.1, 0.78, fill=CARD)
    text(s, 0.88, y + 0.12, 11.5, 0.62,
         [(title, 12.5, True, colour, None, 0), (body, 10.5, False, MUTED, None, 3)])
    y += 0.86

# ---------------------------------------------------------------------------
# 16 — The 3D twin
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "14 · The 3D twin", "What is rendered, and what is only representative.")

left_items = [
    ("Driven by the model", EMERALD),
    ("Beacon colour — normal / warning / critical per sensor status", None),
    ("Label text — live value and unit, updated every render", None),
    ("Selection ring and camera focus — follows the selected sensor", None),
    ("Scenario camera — each scenario frames its own focus sensor", None),
]
right_items = [
    ("Representative geometry only", AMBER),
    ("Arch shell, gorge walls and crest road are modelled, not surveyed", None),
    ("Sensor positions are plausible placements, not as-built coordinates", None),
    ("Internal zoning (grout curtain, gallery, monolith) is illustrative", None),
    ("No structural solver — nothing here is computed from geometry", None),
]
for col, (items, colour) in enumerate([(left_items, EMERALD), (right_items, AMBER)]):
    x = 0.62 + col * 6.25
    card(s, x, 1.5, 6.0, 3.05, fill=CARD, border=colour)
    text(s, x + 0.28, 1.72, 5.45, 0.4, [(items[0][0], 14, True, colour, None, 0)])
    y = 2.25
    for label, _ in items[1:]:
        text(s, x + 0.28, y, 5.45, 0.55, [("•  " + label, 11, False, MUTED, None, 0)])
        y += 0.56

formula(s, 0.62, 4.85, 12.1, [
    "Zoom-aware labels: scale = lerp(1.3, 0.45, smoothstep(distance, 8, 30))",
    "Labels sit at 1.3× when the camera is at its closest dolly stop and shrink to 0.45× at the widest,",
    "so readings stay legible up close without the label swarm covering the dam in the wide shot.",
], accent=CYAN, height=1.35)

footnote(s, "Rendered with react-three-fiber over WebGL. Labels are DOM elements positioned in 3D, not textures.")

# ---------------------------------------------------------------------------
# 17 — Limits
# ---------------------------------------------------------------------------
s = new_slide()
header(s, "15 · What this model is not", "Read this before quoting any number from the demo.", accent=RED)

limits = [
    ("It is not a structural analysis", RED,
     "No finite element method, no seepage solver, no material properties. Sensor responses are linear "
     "coefficients chosen to behave plausibly, not derived from the geometry on screen."),
    ("Superposition is assumed", AMBER,
     "The three load cases add independently. Real coupling — surge changing the seismic response, "
     "seepage lagging behind reservoir level — is not represented. There is no time lag anywhere in the model."),
    ("The risk score is not a standard", AMBER,
     "The 3.8 / 140 / 0.35 weights and the 28 / 60 band edges are demonstration values. They are not "
     "traceable to ICOLD, to any national dam safety guideline, or to a calibrated fragility curve."),
    ("The dam is fictional", PURPLE,
     "Demo Dam is not a real facility. Elevations, thresholds and base readings are invented to be "
     "self-consistent and realistic in shape, nothing more."),
]
y = 1.5
for title, colour, body in limits:
    card(s, 0.62, y, 12.1, 1.15, fill=CARD, border=colour)
    text(s, 0.9, y + 0.16, 11.5, 0.95,
         [(title, 13.5, True, colour, None, 0), (body, 11, False, MUTED, None, 4)])
    y += 1.28

formula(s, 0.62, 6.5, 12.1, [
    "What it does demonstrate: the interaction model, the data flow, and the operator experience — end to end.",
], accent=EMERALD, height=0.62)

# ---------------------------------------------------------------------------
# 18 — Close
# ---------------------------------------------------------------------------
s = new_slide()
card(s, 0.8, 1.3, 11.733, 4.6, fill=CARD, border=CYAN)
text(s, 1.3, 1.75, 10.7, 3.0, [
    ("REPRODUCE ANY FRAME", 12.5, True, MUTED, MONO_FONT, 0),
    ("Three numbers. That's the whole state.", 30, True, CYAN, None, 8),
    ("surge · pga · spillwayGate", 17, False, TEXT, MONO_FONT, 12),
    ("Set those three sliders and every reading, every status colour, every bar fill and the risk "
     "score reproduce exactly. Nothing on the dashboard is stored, streamed, or remembered — it is "
     "all recomputed from those three values on every render.", 12.5, False, MUTED, None, 14),
    ("Source of truth: src/lib/damModel.ts   ·   Deck generator: create_technical_presentation.py",
     11, False, MUTED, MONO_FONT, 16),
])
text(s, 0.8, 6.15, 11.733, 0.6, [
    ("SU YAPI Engineering & Consulting   ·   Para Sidara", 12.5, True, TEXT, None, 0),
    ("Demo Dam (Block 04) — demonstration dataset, not live field telemetry.", 10.5, False, MUTED, None, 4),
], align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Wrote {OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
