import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    BG_COLOR = RGBColor(11, 15, 25)       # #0B0F19 Dark slate
    CARD_BG = RGBColor(22, 29, 45)        # #161D2D Glass card
    CYAN = RGBColor(0, 229, 255)          # #00E5FF Accent Cyan
    PURPLE = RGBColor(139, 92, 246)       # #8B5CF6 Accent Purple
    EMERALD = RGBColor(16, 185, 129)      # #10B981 Success
    AMBER = RGBColor(245, 158, 11)       # #F59E0B Warning
    RED = RGBColor(239, 68, 68)          # #EF4444 Danger
    TEXT_LIGHT = RGBColor(243, 244, 246) # White text
    TEXT_MUTED = RGBColor(156, 163, 175) # Muted gray

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, subtitle_text):
        # Header title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.font.name = 'Arial'

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(13)
            p2.font.color.rgb = TEXT_MUTED
            p2.font.name = 'Arial'
            p2.space_before = Pt(4)

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CYAN):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1)
        else:
            card.line.fill.background()
        return card

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide1)

    # Accent glow box behind title
    add_card(slide1, 0.8, 1.2, 11.733, 5.1, bg_color=CARD_BG, border_color=CYAN)

    txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(10.933), Inches(4.3))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "AQUASHIELD TWIN"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = CYAN

    p2 = tf.add_paragraph()
    p2.text = "Dam Safety Digital Monitoring & Predictive Intelligence"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "A modern WebGL 3D Digital Twin and Real-Time Telemetry Dashboard for Dam Infrastructure"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(20)

    # Metrics pill highlights
    metrics = [
        ("-85%", "Failure Risk", CYAN),
        ("+22 Yrs", "Asset Lifespan", PURPLE),
        ("$1.4M/yr", "OPEX Savings", EMERALD),
        ("< 15s", "Early Warning", AMBER)
    ]
    for i, (val, lbl, col) in enumerate(metrics):
        box_left = 1.2 + i * 2.65
        add_card(slide1, box_left, 4.6, 2.4, 1.2, bg_color=RGBColor(15, 22, 36), border_color=col)
        tb = slide1.shapes.add_textbox(Inches(box_left), Inches(4.7), Inches(2.4), Inches(1.0))
        tf_m = tb.text_frame
        tf_m.word_wrap = True
        p_v = tf_m.paragraphs[0]
        p_v.alignment = PP_ALIGN.CENTER
        p_v.text = val
        p_v.font.size = Pt(22)
        p_v.font.bold = True
        p_v.font.color.rgb = col
        
        p_l = tf_m.add_paragraph()
        p_l.alignment = PP_ALIGN.CENTER
        p_l.text = lbl
        p_l.font.size = Pt(11)
        p_l.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: Executive Summary & What Does This Demo Do?
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide2)
    add_header(slide2, "1. System Overview: What Does This Demo Do?", 
               "AquaShield Twin transforms complex dam safety engineering into interactive 3D WebGL visualizations and predictive AI risk analysis.")

    items_s2 = [
        ("Real-Time 3D Digital Twin Visualizer", 
         "Renders full 3D dam model (Atatürk Earth-fill Dam) with dynamic water levels, interactive spillway gate releases, animated crest displacement, and real-time stress heatmaps."),
        ("Live Telemetry Sensor Network", 
         "Monitors 6 key structural sensor nodes (Piezometers P-01/P-02/P-03, Inclinometer INC-04, Seepage Flow Gauge SF-02, Reservoir Radar WL-01) with live streaming line charts."),
        ("AI Physics 'What-If' Stress Simulator", 
         "Allows operators to adjust Reservoir Surge (0 to +15m), Seismic PGA (0 to 0.45g), and Spillway Gate Openings (10%-100%) to predict instant risk scores."),
        ("Automated Emergency Action Protocol", 
         "Evaluates safety thresholds dynamically and outputs operational recommendations (Normal, Advisory Warning, or Critical Emergency Protocol)."),
        ("Executive ROI & Safety Insights", 
         "Demonstrates how digital twins lower catastrophic risk by 85%, cut annual maintenance costs by $1.4M, and extend dam operational lifespan by 22+ years.")
    ]

    for i, (title, desc) in enumerate(items_s2):
        top = 1.6 + i * 1.08
        add_card(slide2, 0.8, top, 11.733, 0.95, bg_color=CARD_BG, border_color=CYAN if i==0 else None)
        
        tb = slide2.shapes.add_textbox(Inches(1.0), Inches(top + 0.1), Inches(11.333), Inches(0.75))
        tf_i = tb.text_frame
        tf_i.word_wrap = True
        
        p = tf_i.paragraphs[0]
        p.text = f"•  {title}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN if i==0 else TEXT_LIGHT
        
        p2 = tf_i.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(2)

    # =========================================================================
    # SLIDE 3: Key Features Breakdown
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide3)
    add_header(slide3, "2. Key Features & Core Components", 
               "Three integrated modules work together to provide complete structural health oversight.")

    col_width = 3.64
    columns = [
        ("3D WebGL Model Viewport", CYAN, [
            "• Interactive 3D View: 360° orbit, zoom, and pan controls.",
            "• Core Cross-Section: Cutaway view revealing internal clay core and foundation.",
            "• Stress Heatmap Mode: Visual representation of hydrostatic stress zones.",
            "• Dynamic Water & Gates: Animated reservoir surface and spillway particle discharge."
        ]),
        ("Live Telemetry & Sensor Suite", PURPLE, [
            "• Piezometers (P-01, P-02, P-03): Pore water pressure in kPa.",
            "• Crest Inclinometer (INC-04): Structural crest movement in mm.",
            "• Seepage Gauge (SF-02): Drainage gallery leakage flow in L/min.",
            "• Time-Series Graphs: Real-time updated streaming charts."
        ]),
        ("AI Stress Simulator & Risk Controls", EMERALD, [
            "• Reservoir Surge Slider: Simulate extreme flood inflow (+0 to +15m).",
            "• Seismic PGA Slider: Simulate earthquake shaking (0.00g to 0.45g).",
            "• Spillway Control: Manage gate discharge (10% - 100%).",
            "• Dynamic Risk Score: Live 0-100 index with emergency action guidelines."
        ])
    ]

    for i, (col_title, color, bullets) in enumerate(columns):
        left = 0.8 + i * 4.04
        add_card(slide3, left, 1.6, col_width, 5.2, bg_color=CARD_BG, border_color=color)
        
        tb = slide3.shapes.add_textbox(Inches(left + 0.2), Inches(1.8), Inches(col_width - 0.4), Inches(4.8))
        tf_c = tb.text_frame
        tf_c.word_wrap = True
        
        p = tf_c.paragraphs[0]
        p.text = col_title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = color
        
        for b in bullets:
            pb = tf_c.add_paragraph()
            pb.text = b
            pb.font.size = Pt(12)
            pb.font.color.rgb = TEXT_LIGHT
            pb.space_before = Pt(12)

    # =========================================================================
    # SLIDE 4: Interactive 3D Digital Twin Visualizer Modes
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide4)
    add_header(slide4, "3. 3D WebGL Digital Twin Modes", 
               "Switch between 3 distinct visualization perspectives to inspect dam integrity.")

    modes = [
        ("3D Digital Model", CYAN, 
         "Provides a full realistic 3D exterior model of the dam embankment, reservoir water body, intake tower, and spillway structures. Includes live 3D sensor badges attached directly to structural coordinates."),
        ("Internal Core Cross-Section", PURPLE, 
         "Exposes a cutaway cross-section of the dam interior. Allows engineers to inspect the impervious clay core, upstream shell, downstream toe, and internal drainage gallery layout."),
        ("Stress Heatmap Mode", AMBER, 
         "Overlays dynamic color-coded stress gradients (Blue = Stable, Yellow = Moderate Pressure, Red = High Critical Load) based on current reservoir height and seismic ground acceleration.")
    ]

    for i, (m_title, m_color, m_desc) in enumerate(modes):
        top = 1.6 + i * 1.8
        add_card(slide4, 0.8, top, 11.733, 1.5, bg_color=CARD_BG, border_color=m_color)
        
        tb = slide4.shapes.add_textbox(Inches(1.1), Inches(top + 0.2), Inches(11.133), Inches(1.1))
        tf_m = tb.text_frame
        tf_m.word_wrap = True
        
        p = tf_m.paragraphs[0]
        p.text = f"Mode {i+1}: {m_title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = m_color
        
        p2 = tf_m.add_paragraph()
        p2.text = m_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_LIGHT
        p2.space_before = Pt(6)

    # =========================================================================
    # SLIDE 5: Sensor Telemetry & Structural Monitoring
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide5)
    add_header(slide5, "4. Sensor Telemetry & Health Monitoring", 
               "Comprehensive sensor telemetry tracking key engineering parameters in real time.")

    # Sensor Table Card
    add_card(slide5, 0.8, 1.6, 11.733, 5.2, bg_color=CARD_BG, border_color=CYAN)
    
    tb = slide5.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(4.8))
    tf_s = tb.text_frame
    tf_s.word_wrap = True

    p = tf_s.paragraphs[0]
    p.text = "Configured Sensor Inventory & Threshold Envelopes"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN

    sensors_data = [
        ("P-01", "Piezometer", "Upstream Toe (EL 490m)", "242.4 kPa", "300.0 kPa", "Monitors uplift pore pressure under upstream slope."),
        ("P-02", "Piezometer", "Impervious Core (EL 510m)", "185.1 kPa", "220.0 kPa", "Measures internal seepage pressure inside clay core."),
        ("P-03", "Piezometer", "Downstream Core (EL 480m)", "92.6 kPa", "150.0 kPa", "Ensures downstream phreatic line remains safe."),
        ("INC-04", "Inclinometer", "Crest Center (Sec-04)", "3.2 mm", "8.0 mm", "Detects structural lateral deformation of the crest."),
        ("SF-02", "Seepage Flow", "Drainage Gallery G-02", "14.2 L/min", "25.0 L/min", "Measures water collection rate in drainage tunnels."),
        ("WL-01", "Water Radar", "Reservoir Intake Tower", "537.4 m", "545.0 m", "Tracks reservoir surface water elevation.")
    ]

    for s_id, s_type, s_loc, s_val, s_thresh, s_note in sensors_data:
        p_row = tf_s.add_paragraph()
        p_row.text = f"• [{s_id}] {s_type} | {s_loc} — Value: {s_val} (Limit: {s_thresh}) — {s_note}"
        p_row.font.size = Pt(12)
        p_row.font.color.rgb = TEXT_LIGHT
        p_row.space_before = Pt(8)

    # =========================================================================
    # SLIDE 6: AI Physics 'What-If' Simulator
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide6)
    add_header(slide6, "5. AI Physics 'What-If' Stress Simulator", 
               "Test dam stability against severe environmental stressors in real time.")

    add_card(slide6, 0.8, 1.6, 5.6, 5.2, bg_color=CARD_BG, border_color=PURPLE)
    tb_l = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "Simulation Control Knobs"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    sim_controls = [
        ("1. Reservoir Surge (+0m to +15m)", "Simulates extreme rainfall & flood surge. Increases water level, hydrostatic pressure, and seepage rates."),
        ("2. Seismic PGA (0.00g to 0.45g)", "Simulates earthquake ground acceleration. Induces crest vibration, dynamic displacement, and pore pressure rises."),
        ("3. Spillway Gate Opening (10% - 100%)", "Controls reservoir flood discharge rate. Closing gates increases upstream water accumulation.")
    ]

    for title, desc in sim_controls:
        p_c = tf_l.add_paragraph()
        p_c.text = f"• {title}"
        p_c.font.size = Pt(14)
        p_c.font.bold = True
        p_c.font.color.rgb = TEXT_LIGHT
        p_c.space_before = Pt(12)

        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(2)

    # Right side: Dynamic Risk Output
    add_card(slide6, 6.8, 1.6, 5.733, 5.2, bg_color=CARD_BG, border_color=EMERALD)
    tb_r = slide6.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "Dynamic Risk Output & Action Protocols"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    states = [
        ("OPERATIONAL (< 28)", EMERALD, "Structure stable under nominal hydrodynamic loads. Standard monitoring active."),
        ("ELEVATED RISK (28 - 60)", AMBER, "Hydrodynamic pressure elevated. Increase piezometric sampling. Prepare auxiliary spillways."),
        ("CRITICAL ALERT (> 60)", RED, "EMERGENCY PROTOCOL: Open spillway gates to 100%. Dispatch civil team to inspect Downstream Toe (P-01).")
    ]

    for title, color, desc in states:
        p_s = tf_r.add_paragraph()
        p_s.text = f"• Status: {title}"
        p_s.font.size = Pt(14)
        p_s.font.bold = True
        p_s.font.color.rgb = color
        p_s.space_before = Pt(14)

        p_sd = tf_r.add_paragraph()
        p_sd.text = desc
        p_sd.font.size = Pt(11)
        p_sd.font.color.rgb = TEXT_LIGHT
        p_sd.space_before = Pt(2)

    # =========================================================================
    # SLIDE 7: How to Use the Demo (User Step-by-Step Guide)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide7)
    add_header(slide7, "6. How to Use the Demo: Step-by-Step Guide", 
               "Follow these 5 simple steps to explore and evaluate the AquaShield Twin platform.")

    user_steps = [
        ("Step 1: Launch Application", "Run 'npm run dev' and open http://localhost:3000 in any WebGL-enabled desktop browser."),
        ("Step 2: Explore 3D Viewport Modes", "Toggle between '3D Digital Model', 'Internal Core Cross-Section', and 'Stress Heatmap' top buttons."),
        ("Step 3: Interact with 3D Sensors", "Click sensor markers (P-01, INC-04, etc.) in the 3D model or left sidebar to view instant live charts."),
        ("Step 4: Run Stress Scenarios", "Adjust sliders on the right panel (Surge +10m, Seismic PGA 0.35g) to simulate extreme storm & quake events."),
        ("Step 5: Observe Safety Protocol & Reset", "Review how the Risk Score updates automatically to Critical Alert, then click 'Reset Baseline Simulation'.")
    ]

    for i, (stitle, sdesc) in enumerate(user_steps):
        top = 1.6 + i * 1.08
        add_card(slide7, 0.8, top, 11.733, 0.95, bg_color=CARD_BG, border_color=CYAN if i==0 else None)
        
        tb = slide7.shapes.add_textbox(Inches(1.0), Inches(top + 0.1), Inches(11.333), Inches(0.75))
        tf_us = tb.text_frame
        tf_us.word_wrap = True
        
        p = tf_us.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN if i==0 else TEXT_LIGHT
        
        p2 = tf_us.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(2)

    # =========================================================================
    # SLIDE 8: Technical Architecture & Conclusion
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_bg(slide8)
    add_header(slide8, "7. Technology Stack & Business Value", 
               "Built with modern web technologies to deliver high-performance digital twin visualization.")

    add_card(slide8, 0.8, 1.6, 5.6, 5.2, bg_color=CARD_BG, border_color=CYAN)
    tb_t = slide8.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True

    p = tf_t.paragraphs[0]
    p.text = "Technology Stack"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN

    tech_stack = [
        ("Next.js 16 (React 19, Turbopack)", "App Router, SSR-disabled dynamic client rendering for WebGL."),
        ("Three.js & React Three Fiber", "High-performance GPU WebGL rendering, custom dam geometries, particle water physics."),
        ("Chart.js & react-chartjs-2", "Streaming real-time telemetry line charts with smooth animated transitions."),
        ("Tailwind CSS v4", "Glassmorphism UI design tokens, dark tech aesthetic, responsive executive grid layout.")
    ]

    for title, desc in tech_stack:
        p_t = tf_t.add_paragraph()
        p_t.text = f"• {title}"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_LIGHT
        p_t.space_before = Pt(10)

        p_td = tf_t.add_paragraph()
        p_td.text = desc
        p_td.font.size = Pt(11)
        p_td.font.color.rgb = TEXT_MUTED
        p_td.space_before = Pt(2)

    add_card(slide8, 6.8, 1.6, 5.733, 5.2, bg_color=CARD_BG, border_color=PURPLE)
    tb_b = slide8.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True

    p = tf_b.paragraphs[0]
    p.text = "Business Impact & Value Delivered"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    biz_value = [
        ("Risk Reduction", "-85% reduction in catastrophic failure risk through early anomaly detection."),
        ("Asset Longevity", "+22 years extended operational lifespan via proactive structural care."),
        ("Cost Efficiency", "$1.4M annual savings in operational maintenance and inspection OPEX."),
        ("Response Speed", "<15 second automated emergency warning dispatch during extreme events.")
    ]

    for title, desc in biz_value:
        p_b = tf_b.add_paragraph()
        p_b.text = f"• {title}"
        p_b.font.size = Pt(13)
        p_b.font.bold = True
        p_b.font.color.rgb = TEXT_LIGHT
        p_b.space_before = Pt(10)

        p_bd = tf_b.add_paragraph()
        p_bd.text = desc
        p_bd.font.size = Pt(11)
        p_bd.font.color.rgb = TEXT_MUTED
        p_bd.space_before = Pt(2)

    # Save output
    output_filename = "AquaShield_Twin_Dam_Monitoring_Presentation.pptx"
    prs.save(output_filename)
    print(f"Successfully generated PowerPoint presentation: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    create_deck()
